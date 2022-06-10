# -*- coding: utf-8 -*-
'''
*******************************************************************************
 * @fileName : get_text_from_pptx.py
 * @author   : "Ko Sun Ho"
 * @comment  : 문서 파싱 pptx module (lib - python-pptx)
 * @revision history
 * date            author         comment
 * ----------      ---------      ----------------------------------------------
 * 2022. 05. 16.  Ko Sun Ho       작성
 *******************************************************************************
'''
# import
from pptx import Presentation


class GetPptxText:
    def __init__(self, filepath, filename):
        self.filepath = filepath
        self.filename = filename
        self.fullpath = self.filepath + '/' + self.filename + '.pptx'
        self.f = Presentation(self.fullpath)

    def _get_word_text(self):

        result = []
        # 슬라이드 순회하며 텍스트 파싱
        for slide in self.f.slides:
            print(slide.shapes)
            for shape in slide.shapes:
                # 슬라이드 text 존재 확인
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text == '':
                        continue
                    result.append(paragraph.text)

        return result

    def get_text_pptx_main_str(self):
        return '\n'.join(self._get_word_text())

    def get_text_pptx_main_list(self):
        return self._get_word_text()

    def get_text_pptx_main_with_filename(self):
        return [self.filename, '\n'.join(self._get_word_text())]


if __name__ == '__main__':
    from doc_miner.common_modules.common import *

    # resources 경로가져오기
    fm = 'pptx'
    path = os.getcwd()
    rsc_path = OperatingSystem.get_parent_path(path)
    sample_path = os.path.join(rsc_path, 'resources/sample_doc')
    docx_files = OperatingSystem.sperate_file_extension(check_dir=sample_path)[fm]

    # .docx text 파싱
    sample_docx_filename = docx_files[0]
    ins_pdf = GetPptxText(filepath=sample_path, filename=sample_docx_filename)

    res_list = ins_pdf.get_text_pptx_main_list()
    print(res_list)

    res_str = ins_pdf.get_text_pptx_main_str()
    print(res_str)
